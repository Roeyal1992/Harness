#!/usr/bin/env python3
"""
Deloitte PPTX Logo Utility

Fetches corporate logos and places them on slides. Tries multiple free
sources, then falls back to generating a branded initial-badge so the
slide always has something usable.

Usage as a module:
    from fetch_logo import add_logo_to_slide, fetch_logo_image

    # Fetch and place in one call
    add_logo_to_slide(slide, "Salesforce", "salesforce.com", x=1, y=2, h=0.6)

    # Or fetch first, place later
    img_path = fetch_logo_image("Apple", "apple.com", output_dir="/tmp")
    slide.shapes.add_image(img_path, Inches(1), Inches(2), height=Inches(0.6))

Usage from CLI:
    python fetch_logo.py salesforce.com
    python fetch_logo.py apple.com google.com microsoft.com --size 256
"""

import os
import sys
import hashlib
import argparse
import tempfile
from pathlib import Path

try:
    from pptx.util import Inches
except ImportError:
    pass

# ── Configuration ────────────────────────────────────────────────────────

LOGO_CACHE_DIR = os.path.join(tempfile.gettempdir(), 'deloitte_pptx_logos')

# Free logo APIs, tried in order. {domain} is replaced at runtime.
LOGO_SOURCES = [
    'https://logo.clearbit.com/{domain}?size=256',           # Best quality, no key
    'https://img.logo.dev/{domain}?token=pk_anonymous',      # Backup
    'https://www.google.com/s2/favicons?domain={domain}&sz=128',  # Low-res fallback
]

# Brand colors for generated fallback badges (Deloitte palette)
BADGE_COLORS = [
    (0x86, 0xBC, 0x25),  # Deloitte Green
    (0x0D, 0x83, 0x90),  # Teal
    (0x00, 0xA3, 0xE0),  # Blue
    (0x26, 0x89, 0x0D),  # Mid Green
    (0x00, 0x76, 0x80),  # Teal 6
    (0x00, 0x55, 0x87),  # Blue 5
    (0x04, 0x6A, 0x38),  # Deep Green
]


# ── Core Functions ───────────────────────────────────────────────────────

def fetch_logo_image(company_name, domain, output_dir=None, size=256):
    """Fetch a company logo, returning the local file path.

    Tries free logo APIs in order, caches results, and falls back to
    generating a colored initial-badge if all network sources fail.

    Args:
        company_name: Display name (e.g. "Salesforce") — used for fallback badge
        domain: Company domain (e.g. "salesforce.com") — used for API lookup
        output_dir: Where to save. Defaults to system temp dir.
        size: Requested image size in pixels (APIs may return different sizes)

    Returns:
        str: Absolute path to the logo image file (PNG)
    """
    if output_dir is None:
        output_dir = LOGO_CACHE_DIR
    os.makedirs(output_dir, exist_ok=True)

    # Normalize domain
    domain = domain.lower().strip()
    if domain.startswith('http'):
        from urllib.parse import urlparse
        domain = urlparse(domain).netloc or domain
    domain = domain.replace('www.', '')

    # Cache key
    cache_file = os.path.join(output_dir, f"logo_{domain.replace('.', '_')}.png")
    if os.path.exists(cache_file) and os.path.getsize(cache_file) > 500:
        return cache_file

    # Try each API source
    img_data = None
    for url_template in LOGO_SOURCES:
        url = url_template.format(domain=domain)
        img_data = _download(url)
        if img_data and len(img_data) > 500:  # Valid image (not error page)
            break
        img_data = None

    if img_data:
        with open(cache_file, 'wb') as f:
            f.write(img_data)
        return cache_file

    # All sources failed — generate fallback badge
    fallback_path = _generate_badge(company_name, domain, output_dir, size)
    return fallback_path


def add_logo_to_slide(slide, company_name, domain, x, y, h=0.6,
                      output_dir=None):
    """Fetch a logo and place it on a slide.

    Args:
        slide: pptx Slide object
        company_name: e.g. "Microsoft"
        domain: e.g. "microsoft.com"
        x, y: Position in inches
        h: Height in inches (width auto-calculated to preserve aspect ratio)
        output_dir: Logo cache directory

    Returns:
        The added picture shape, or a fallback shape if image fails
    """
    logo_path = fetch_logo_image(company_name, domain, output_dir)

    try:
        pic = slide.shapes.add_picture(
            logo_path,
            Inches(x), Inches(y),
            height=Inches(h)
        )
        return pic
    except Exception as e:
        # If even the fallback image fails, add a text placeholder
        print(f"  Warning: Could not place logo for {company_name}: {e}")
        return _add_text_badge(slide, company_name, x, y, h)


def add_logos_row(slide, companies, y=2.0, h=0.5, x_start=0.5, x_end=12.83,
                  output_dir=None):
    """Place a row of company logos evenly spaced across the slide.

    Args:
        companies: List of (company_name, domain) tuples
        y: Vertical position in inches
        h: Logo height in inches
        x_start, x_end: Horizontal bounds

    Returns:
        List of added shapes
    """
    n = len(companies)
    if n == 0:
        return []

    spacing = (x_end - x_start) / n
    shapes = []
    for i, (name, domain) in enumerate(companies):
        x = x_start + i * spacing + (spacing - h) / 2  # Center in slot
        shape = add_logo_to_slide(slide, name, domain, x, y, h, output_dir)
        shapes.append(shape)
    return shapes


def add_logo_grid(slide, companies, cols=4, x_start=0.5, y_start=2.0,
                  cell_w=3.0, cell_h=1.5, logo_h=0.5, output_dir=None):
    """Place logos in a grid with company names below each.

    Args:
        companies: List of (company_name, domain) tuples
        cols: Number of columns
        x_start, y_start: Top-left of grid
        cell_w, cell_h: Size of each grid cell
        logo_h: Height of each logo image

    Returns:
        List of added shapes
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    shapes = []
    for i, (name, domain) in enumerate(companies):
        col = i % cols
        row = i // cols
        cx = x_start + col * cell_w
        cy = y_start + row * cell_h

        # Center logo in cell
        logo_x = cx + (cell_w - logo_h) / 2
        shape = add_logo_to_slide(slide, name, domain, logo_x, cy, logo_h, output_dir)
        shapes.append(shape)

        # Company name label below logo
        tb = slide.shapes.add_textbox(
            Inches(cx), Inches(cy + logo_h + 0.1),
            Inches(cell_w), Inches(0.3)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        run.font.size = Pt(9)
        run.font.name = 'Aptos'
        run.font.color.rgb = RGBColor(0x53, 0x56, 0x5A)
        shapes.append(tb)

    return shapes


# ── Internal Helpers ─────────────────────────────────────────────────────

def _download(url, timeout=8):
    """Download bytes from a URL. Returns None on failure."""
    try:
        import urllib.request
        req = urllib.request.Request(url, headers={
            'User-Agent': 'Mozilla/5.0 (Deloitte PPTX Skill)',
            'Accept': 'image/png,image/*,*/*',
        })
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            if resp.status == 200:
                data = resp.read()
                content_type = resp.headers.get('Content-Type', '')
                # Verify we got an image, not an error page
                if 'image' in content_type or data[:4] == b'\x89PNG' or data[:2] == b'\xff\xd8':
                    return data
    except Exception:
        pass
    return None


def _generate_badge(company_name, domain, output_dir, size=256):
    """Generate a colored circle badge with company initials.

    Uses Pillow if available, otherwise creates a minimal 1-color PNG.
    """
    # Pick a consistent color based on the domain hash
    color_idx = int(hashlib.md5(domain.encode()).hexdigest(), 16) % len(BADGE_COLORS)
    bg_color = BADGE_COLORS[color_idx]

    # Extract initials (up to 2 characters)
    words = company_name.replace('.', ' ').replace('-', ' ').split()
    if len(words) >= 2:
        initials = (words[0][0] + words[1][0]).upper()
    elif words:
        initials = words[0][:2].upper()
    else:
        initials = domain[0].upper()

    cache_file = os.path.join(output_dir, f"badge_{domain.replace('.', '_')}.png")

    try:
        from PIL import Image, ImageDraw, ImageFont

        img = Image.new('RGBA', (size, size), (0, 0, 0, 0))
        draw = ImageDraw.Draw(img)

        # Draw circle
        margin = int(size * 0.05)
        draw.ellipse([margin, margin, size - margin, size - margin],
                     fill=bg_color)

        # Draw initials
        font_size = int(size * 0.38)
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
                                       font_size)
        except:
            try:
                font = ImageFont.truetype("arial.ttf", font_size)
            except:
                font = ImageFont.load_default()

        bbox = draw.textbbox((0, 0), initials, font=font)
        tw = bbox[2] - bbox[0]
        th = bbox[3] - bbox[1]
        tx = (size - tw) / 2
        ty = (size - th) / 2 - bbox[1]
        draw.text((tx, ty), initials, fill=(255, 255, 255), font=font)

        img.save(cache_file, 'PNG')

    except ImportError:
        # No Pillow — create a minimal solid-color PNG
        _create_minimal_png(cache_file, bg_color, size)

    return cache_file


def _create_minimal_png(path, rgb_tuple, size=64):
    """Create a tiny solid-color PNG without any dependencies."""
    import struct
    import zlib

    width = height = size
    r, g, b = rgb_tuple

    # Build raw RGBA pixel data
    raw = b''
    for _ in range(height):
        raw += b'\x00'  # filter byte
        for _ in range(width):
            raw += bytes([r, g, b, 255])

    def chunk(chunk_type, data):
        c = chunk_type + data
        crc = struct.pack('>I', zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack('>I', len(data)) + c + crc

    with open(path, 'wb') as f:
        f.write(b'\x89PNG\r\n\x1a\n')  # PNG signature
        f.write(chunk(b'IHDR', struct.pack('>IIBBBBB', width, height, 8, 6, 0, 0, 0)))
        f.write(chunk(b'IDAT', zlib.compress(raw)))
        f.write(chunk(b'IEND', b''))


def _add_text_badge(slide, company_name, x, y, h):
    """Last-resort fallback: add a text box with the company name."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    # Colored rectangle with company name
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(h * 2.5), Inches(h)
    )
    shape.fill.solid()
    idx = hash(company_name) % len(BADGE_COLORS)
    r, g, b = BADGE_COLORS[idx]
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = company_name
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = 'Aptos'

    return shape


# ── CLI ──────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description='Fetch corporate logos for PPTX')
    parser.add_argument('domains', nargs='+', help='Company domains (e.g. apple.com)')
    parser.add_argument('--size', type=int, default=256, help='Image size in px')
    parser.add_argument('--output', '-o', default=LOGO_CACHE_DIR, help='Output directory')
    args = parser.parse_args()

    os.makedirs(args.output, exist_ok=True)

    for domain in args.domains:
        name = domain.split('.')[0].title()
        path = fetch_logo_image(name, domain, args.output, args.size)
        print(f"  {domain} → {path}")


if __name__ == '__main__':
    main()
