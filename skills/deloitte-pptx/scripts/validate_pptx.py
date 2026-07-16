#!/usr/bin/env python3
"""
Deloitte PPTX Validator
Checks a generated presentation for brand compliance, formatting errors,
placeholder leftovers, color violations, font issues, and layout problems.

Usage:
    python validate_pptx.py output.pptx [--visual] [--strict]

Exit codes:
    0 = PASS (no errors, warnings only)
    1 = FAIL (errors found)
"""

import sys
import os
import re
import json
import argparse
from collections import defaultdict

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx --break-system-packages")
    sys.exit(1)

# ── Brand Constants ──────────────────────────────────────────────────────────

APPROVED_COLORS_HEX = {
    # Primary
    'FFFFFF', '000000',
    # Greens
    '86BC25', '26890D', '046A38', '1C3D26', '0DF200', 'F1F6E4',
    '009A44', '43B02A', 'C4D600', 'E3E48D',
    # Teals
    '0D8390', '004F59', '007680', '0097A9', '00ABAB', '6FC2B4', '9DD4CF', 'DDEFE8',
    # Blues
    '007CB0', '041E42', '012169', '005587', '0076A8', '00A3E0', '62B5E5', 'A0DCFF',
    # Grays
    '222222', '282728', '53565A', '63666A', '75787B', '97999B',
    'A7A8AA', 'BBBCBC', 'D0D0CE', 'E6E6E6',
    # Brights
    '3EFAC5', '33F0FF', '86EB22',
    # Functional (limited use)
    'DA291C', 'DA2910', 'ED8B00', 'FFCD00',
    # Gradient accents
    'B7E320', '63C631', '6BC930', 'FCFF5D',
}

# Close-enough tolerance: allow colors within distance 5 of an approved color
# (catches minor rounding from theme/template inheritance)
COLOR_TOLERANCE = 8

APPROVED_FONTS = {
    'Aptos', 'Aptos Display',
    # Legacy fonts still present in some templates
    'Open Sans', 'Open Sans Light',
    'Calibri', 'Calibri Light',
    # System fallbacks
    'Arial', 'Helvetica',
    # Accent
    'Stix Two Text',
    # PowerPoint theme font references (resolve to template's actual font)
    '+mj-lt', '+mn-lt', '+mj-ea', '+mn-ea', '+mj-cs', '+mn-cs',
}

PLACEHOLDER_PATTERNS = [
    r'xxxx',
    r'lorem\s+ipsum',
    r'dummy\s+text',
    r'text\s+here',
    r'text\s+runs\s+here',
    r'insert\s+text',
    r'click\s+to\s+add',
    r'replace\s+with',
    r'placeholder',
    r'sample\s+text',
    r'your\s+text',
    r'type\s+here',
    r'edit\s+this',
    r'name\s+surname',
    r'xxx@email',
    r'\+49\s*\[xx',
    r'this\s+is\s+dummy',
    r'not\s+here\s+to\s+be\s+read',
    r'this\s+is\s+just\s+text\s+to\s+show',
    r'could\s+insert\s+text',
    r'populated\s+with\s+real\s+text',
    # Deloitte template boilerplate (footers, master-inherited text)
    r'insert\s+appropriate\s+copyright',
    r'to\s+edit.*click\s+view',
    r'slide\s+master\s+>\s+slide\s+master',
    r'presentation\s+title',
    r'member\s+firms\s+and\s+DTTL',
    r'switch\s+the\s+layout\s+of\s+this\s+slide',
    r'right.click\s+on\s+its\s+thumbnail',
    r'predefined\s+text\s+levels',
    r'increase\s+or\s+decrease\s+list\s+level',
]
PLACEHOLDER_RE = re.compile('|'.join(PLACEHOLDER_PATTERNS), re.IGNORECASE)

# Footer zone: content below this y-position encroaches on source/footer area
FOOTER_ZONE_Y = 6.4  # inches — source citations go at ~6.7", footers at ~6.96"

SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.50
MARGIN_IN = 0.3  # Minimum safe margin from edges

# ── Utility ──────────────────────────────────────────────────────────────────

def emu_to_in(emu):
    return emu / 914400

def rgb_to_hex(rgb):
    if rgb is None:
        return None
    return f"{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"

def color_distance(hex1, hex2):
    """Simple RGB distance between two hex color strings."""
    r1, g1, b1 = int(hex1[0:2], 16), int(hex1[2:4], 16), int(hex1[4:6], 16)
    r2, g2, b2 = int(hex2[0:2], 16), int(hex2[2:4], 16), int(hex2[4:6], 16)
    return max(abs(r1-r2), abs(g1-g2), abs(b1-b2))

def is_approved_color(hex_color):
    """Check if a color is in the approved palette (with tolerance)."""
    if hex_color is None:
        return True
    hx = hex_color.upper()
    if hx in APPROVED_COLORS_HEX:
        return True
    # Check tolerance
    for approved in APPROVED_COLORS_HEX:
        if color_distance(hx, approved) <= COLOR_TOLERANCE:
            return True
    return False

# ── Issue Tracking ───────────────────────────────────────────────────────────

class Issue:
    def __init__(self, level, slide_num, category, message, detail=''):
        self.level = level       # 'ERROR', 'WARNING', 'INFO'
        self.slide = slide_num
        self.category = category
        self.message = message
        self.detail = detail

    def __str__(self):
        prefix = {'ERROR': '❌', 'WARNING': '⚠️', 'INFO': 'ℹ️'}[self.level]
        loc = f"Slide {self.slide}" if self.slide else "Deck"
        s = f"  {prefix} [{loc}] {self.category}: {self.message}"
        if self.detail:
            s += f"\n      → {self.detail}"
        return s

class Validator:
    def __init__(self, pptx_path, strict=False):
        self.path = pptx_path
        self.strict = strict
        self.prs = Presentation(pptx_path)
        self.issues = []
        self.stats = defaultdict(int)

    def add(self, level, slide, category, message, detail=''):
        self.issues.append(Issue(level, slide, category, message, detail))
        self.stats[level] += 1

    # ── Checks ───────────────────────────────────────────────────────────

    def check_all(self):
        """Run all validation checks."""
        self.check_deck_structure()
        for i, slide in enumerate(self.prs.slides):
            sn = i + 1
            self.check_placeholder_text(slide, sn)
            self.check_fonts(slide, sn)
            self.check_colors(slide, sn)
            self.check_bounds(slide, sn)
            self.check_overlaps(slide, sn)
            self.check_footer_zone(slide, sn)
            self.check_empty_shapes(slide, sn)
            self.check_text_overflow(slide, sn)
            self.check_narrow_wrapping(slide, sn)
            self.check_font_legibility(slide, sn)
            self.check_zero_padded_numbers(slide, sn)
            self.check_card_whitespace(slide, sn)
            self.check_text_density(slide, sn)
            self.check_font_size_jumps(slide, sn)
            self.check_contrast(slide, sn)
        self.check_consistency()
        self.check_divider_consistency()
        self.check_slide_density()
        self.check_content_quality()
        return self

    def check_deck_structure(self):
        """Check overall deck structure."""
        num_slides = len(self.prs.slides)
        if num_slides == 0:
            self.add('ERROR', None, 'Structure', 'Presentation has no slides')
            return

        self.add('INFO', None, 'Structure', f'Deck has {num_slides} slides')

        # Check dimensions
        w = emu_to_in(self.prs.slide_width)
        h = emu_to_in(self.prs.slide_height)
        if abs(w - 13.33) > 0.1 or abs(h - 7.50) > 0.1:
            self.add('WARNING', None, 'Structure',
                     f'Non-standard dimensions: {w:.2f}" x {h:.2f}"',
                     'Expected 13.33" x 7.50" for 16:9 Deloitte template')

    def check_placeholder_text(self, slide, sn):
        """Check for leftover placeholder/template text."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text
            if not full_text.strip():
                continue

            match = PLACEHOLDER_RE.search(full_text)
            if match:
                snippet = full_text[:80].replace('\n', ' ').strip()
                self.add('ERROR', sn, 'Placeholder',
                         f'Leftover placeholder text found',
                         f'"{snippet}..." (matched: "{match.group()}")')

    def check_fonts(self, slide, sn):
        """Check all text uses approved fonts."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font_name = run.font.name
                    if font_name and font_name not in APPROVED_FONTS:
                        snippet = run.text[:40].strip()
                        if snippet:
                            self.add('WARNING' if not self.strict else 'ERROR',
                                     sn, 'Font',
                                     f'Non-standard font: "{font_name}"',
                                     f'Text: "{snippet}..."')

    def check_colors(self, slide, sn):
        """Check text and shape colors against brand palette."""
        checked = set()

        for shape in slide.shapes:
            # Check shape fills
            if hasattr(shape, 'fill'):
                try:
                    fill = shape.fill
                    if fill.type is not None and fill.fore_color and fill.fore_color.type is not None:
                        hx = rgb_to_hex(fill.fore_color.rgb)
                        if hx and not is_approved_color(hx):
                            key = f"fill-{sn}-{hx}"
                            if key not in checked:
                                checked.add(key)
                                self.add('WARNING', sn, 'Color',
                                         f'Off-palette fill color: #{hx}',
                                         f'Shape: {shape.name}')
                except:
                    pass

            # Check text colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                hx = rgb_to_hex(run.font.color.rgb)
                                if hx and not is_approved_color(hx):
                                    key = f"text-{sn}-{hx}"
                                    if key not in checked:
                                        checked.add(key)
                                        snippet = run.text[:30].strip()
                                        self.add('WARNING', sn, 'Color',
                                                 f'Off-palette text color: #{hx}',
                                                 f'Text: "{snippet}"')
                        except:
                            pass

    def check_bounds(self, slide, sn):
        """Check shapes aren't overflowing slide boundaries."""
        slide_w = emu_to_in(self.prs.slide_width)
        slide_h = emu_to_in(self.prs.slide_height)

        for shape in slide.shapes:
            try:
                left = emu_to_in(shape.left)
                top = emu_to_in(shape.top)
                width = emu_to_in(shape.width)
                height = emu_to_in(shape.height)
                right = left + width
                bottom = top + height
            except:
                continue

            if left < -0.1:
                self.add('WARNING', sn, 'Bounds',
                         f'Shape extends past left edge',
                         f'"{shape.name}" at x={left:.2f}"')
            if top < -0.1:
                self.add('WARNING', sn, 'Bounds',
                         f'Shape extends past top edge',
                         f'"{shape.name}" at y={top:.2f}"')
            if right > slide_w + 0.1:
                self.add('WARNING', sn, 'Bounds',
                         f'Shape extends past right edge',
                         f'"{shape.name}" right edge at {right:.2f}" (slide width: {slide_w:.2f}")')
            if bottom > slide_h + 0.1:
                self.add('WARNING', sn, 'Bounds',
                         f'Shape extends past bottom edge',
                         f'"{shape.name}" bottom at {bottom:.2f}" (slide height: {slide_h:.2f}")')

    def check_overlaps(self, slide, sn):
        """Detect overlapping shapes, especially text-on-text collisions.

        This catches the common visual bug where programmatically placed
        elements share the same vertical/horizontal space, creating unreadable
        pile-ups (e.g. a progress bar label sitting on top of a text paragraph,
        or a shape partially obscuring a line of text below it).

        Strategy:
        - Collect bounding boxes for all shapes with text or visible fills.
        - Compare all pairs for intersection area.
        - Classify by severity:
            TEXT-on-TEXT overlap: ERROR (always unreadable)
            TEXT-on-SHAPE overlap: WARNING (text may be obscured)
            SHAPE-on-SHAPE overlap: ignored if small (intentional accents)
        - Use tolerances to filter out intentional layering patterns:
            - Thin accent bars (h < 0.1") sitting on a card: skip
            - Shapes fully contained in a larger one: treat as intentional
              (e.g. text box inside a rectangle background)
            - Tiny overlap area (< 0.05 sq-in): skip
        """
        # Gather shape info
        shape_info = []
        for shape in slide.shapes:
            try:
                left = emu_to_in(shape.left)
                top = emu_to_in(shape.top)
                width = emu_to_in(shape.width)
                height = emu_to_in(shape.height)
            except:
                continue

            # Skip shapes with zero area
            if width < 0.02 or height < 0.02:
                continue

            # Skip template placeholders — their positions are set by the
            # slide master and are intentionally layered (e.g., header text
            # over subtitle text in the same layout zone)
            is_placeholder = False
            try:
                pf = shape.placeholder_format
                if pf is not None:
                    is_placeholder = True
            except (ValueError, AttributeError):
                pass
            if is_placeholder:
                continue

            has_text = False
            text_snippet = ''
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    has_text = True
                    text_snippet = txt[:35].replace('\n', ' ')

            has_fill = False
            try:
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    has_fill = True
            except:
                pass

            # Skip invisible shapes (no text, no fill)
            if not has_text and not has_fill:
                continue

            # Is this a thin accent/line? (height or width < 0.1")
            is_accent = (height < 0.12 or width < 0.12)

            shape_info.append({
                'name': shape.name,
                'left': left,
                'top': top,
                'right': left + width,
                'bottom': top + height,
                'width': width,
                'height': height,
                'area': width * height,
                'has_text': has_text,
                'text': text_snippet,
                'is_accent': is_accent,
                'shape_id': id(shape),
            })

        # Compare all pairs
        reported = set()
        n = len(shape_info)
        for i in range(n):
            for j in range(i + 1, n):
                a = shape_info[i]
                b = shape_info[j]

                # Calculate intersection
                ix_left = max(a['left'], b['left'])
                ix_top = max(a['top'], b['top'])
                ix_right = min(a['right'], b['right'])
                ix_bottom = min(a['bottom'], b['bottom'])

                if ix_left >= ix_right or ix_top >= ix_bottom:
                    continue  # No intersection

                ix_w = ix_right - ix_left
                ix_h = ix_bottom - ix_top
                ix_area = ix_w * ix_h

                # ── Filter out intentional patterns ──

                # Skip tiny overlaps (< 0.15 sq-in)
                if ix_area < 0.15:
                    continue

                # Skip if one is an accent bar/line (thin decorative element)
                if a['is_accent'] or b['is_accent']:
                    continue

                # Skip if one shape fully contains the other AND the outer
                # one has no text (intentional layering: background rect
                # with text box on top). If BOTH have text, a containment
                # is still a collision — the text will render on top of each other.
                a_contains_b = (a['left'] <= b['left'] + 0.05 and
                                a['top'] <= b['top'] + 0.05 and
                                a['right'] >= b['right'] - 0.05 and
                                a['bottom'] >= b['bottom'] - 0.05)
                b_contains_a = (b['left'] <= a['left'] + 0.05 and
                                b['top'] <= a['top'] + 0.05 and
                                b['right'] >= a['right'] - 0.05 and
                                b['bottom'] >= a['bottom'] - 0.05)

                if a_contains_b and not a['has_text']:
                    continue  # a is a background for b
                if b_contains_a and not b['has_text']:
                    continue  # b is a background for a

                # ── Card pattern detection ──
                # Recognize text boxes placed on filled rectangles as 
                # intentional card layouts, even when the text box slightly
                # overflows the rectangle boundary (common with dynamic
                # content in numbered cards, comparison panels, etc.)
                # A card pattern is: text shape mostly inside a non-text
                # filled shape (>80% overlap), with similar x-position
                # and width, where the non-text shape is at least as wide.
                if (a['has_text'] != b['has_text']):  # one text, one not
                    text_s = a if a['has_text'] else b
                    bg_s = b if a['has_text'] else a
                    # Check if text box is horizontally aligned with bg
                    x_aligned = (abs(text_s['left'] - bg_s['left']) < 0.3 and
                                 abs(text_s['right'] - bg_s['right']) < 0.3)
                    # Check if text box top is at or below bg top
                    y_aligned = text_s['top'] >= bg_s['top'] - 0.1
                    # Check high overlap percentage of text box area
                    text_area = text_s['area']
                    text_overlap_pct = (ix_area / text_area * 100) if text_area > 0 else 0
                    if x_aligned and y_aligned and text_overlap_pct > 75:
                        continue  # intentional card pattern

                # If BOTH have text, don't skip — that's a real collision

                # Skip if overlap is just edge-touching (< 0.2" in either dim)
                if ix_w < 0.2 or ix_h < 0.2:
                    continue

                # ── Classify severity ──

                # Avoid duplicate reports for the same pair
                pair_key = tuple(sorted([a['shape_id'], b['shape_id']]))
                if pair_key in reported:
                    continue
                reported.add(pair_key)

                both_text = a['has_text'] and b['has_text']

                # Calculate what % of the smaller shape is covered
                smaller_area = min(a['area'], b['area'])
                overlap_pct = (ix_area / smaller_area * 100) if smaller_area > 0 else 0

                if both_text:
                    # Text-on-text: always bad
                    level = 'ERROR'
                    category = 'Overlap'
                    a_desc = f'"{a["name"]}"' + (f' ("{a["text"]}...")' if a['text'] else '')
                    b_desc = f'"{b["name"]}"' + (f' ("{b["text"]}...")' if b['text'] else '')
                    self.add(level, sn, category,
                             f'Text elements overlap ({ix_w:.1f}" × {ix_h:.1f}" intersection, '
                             f'{overlap_pct:.0f}% of smaller shape)',
                             f'{a_desc} collides with {b_desc} — '
                             f'text will be unreadable in the overlap zone. '
                             f'Move one element or reduce content to eliminate the collision')

                elif a['has_text'] or b['has_text']:
                    # Text overlapping a filled shape (text may be hidden)
                    text_shape = a if a['has_text'] else b
                    other_shape = b if a['has_text'] else a
                    t_desc = f'"{text_shape["name"]}"'
                    if text_shape['text']:
                        t_desc += f' ("{text_shape["text"]}...")'

                    # Only warn if significant overlap
                    if overlap_pct > 20:
                        self.add('WARNING', sn, 'Overlap',
                                 f'Text element partially obscured by shape '
                                 f'({overlap_pct:.0f}% covered)',
                                 f'{t_desc} overlaps with "{other_shape["name"]}" — '
                                 f'check z-order or reposition to keep text visible')

    def check_footer_zone(self, slide, sn):
        """Flag non-footnote content that intrudes into the footer zone.

        The area below y=6.4" is reserved for source citations (y≈6.7") and
        template footers (y≈6.96"). Content shapes (charts, bars, cards) that
        extend into this zone will collide with footnotes or get clipped.
        Small text boxes (likely source lines) are exempt.
        """
        for shape in slide.shapes:
            try:
                top = emu_to_in(shape.top)
                height = emu_to_in(shape.height)
                width = emu_to_in(shape.width)
                bottom = top + height
            except:
                continue

            # Only flag shapes that START above the footer zone but EXTEND into it
            # (shapes that start in the zone are likely intentional footnotes)
            if top >= FOOTER_ZONE_Y:
                continue

            if bottom <= FOOTER_ZONE_Y:
                continue

            # Skip small text boxes (likely source citations placed intentionally)
            if shape.has_text_frame and height < 0.5 and width < 13.0:
                continue

            # Skip template placeholders
            try:
                pf = shape.placeholder_format
                if pf is not None:
                    continue
            except (ValueError, AttributeError):
                pass

            # Skip thin accent lines
            if height < 0.12 or width < 0.12:
                continue

            intrusion = bottom - FOOTER_ZONE_Y
            self.add('WARNING', sn, 'Footer zone',
                     f'Content extends {intrusion:.1f}" into footer zone '
                     f'(bottom at y={bottom:.1f}", limit is {FOOTER_ZONE_Y}")',
                     f'"{shape.name}" — reduce content height, shrink font, '
                     f'or split across slides to keep clear of source/footer area')

    def check_empty_shapes(self, slide, sn):
        """Flag visible shapes with no content (possible missed placeholders)."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                try:
                    pf = shape.placeholder_format
                    if pf is not None:
                        ptype = pf.type
                        # Only flag TITLE and BODY placeholders, not PICTURE
                        if ptype in (1, 2, 3, 7):  # TITLE, BODY, CENTER_TITLE, OBJECT
                            self.add('WARNING', sn, 'Empty',
                                     f'Empty text placeholder',
                                     f'"{shape.name}" (type={ptype})')
                except (ValueError, AttributeError):
                    pass

    def check_text_overflow(self, slide, sn):
        """Heuristic check for text that likely overflows its container."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                box_w_in = emu_to_in(shape.width)
                box_h_in = emu_to_in(shape.height)
            except:
                continue

            tf = shape.text_frame
            total_text = tf.text.strip()
            if not total_text:
                continue

            font_size = _get_dominant_font_size(tf)
            est = _estimate_wrap(total_text, box_w_in, box_h_in, font_size)

            if est['overflows']:
                self.add('WARNING', sn, 'Overflow',
                         f'Text likely overflows its container',
                         f'"{shape.name}": ~{est["est_lines"]} lines at {font_size}pt, '
                         f'box fits ~{est["capacity"]} '
                         f'({box_w_in:.1f}" × {box_h_in:.1f}")')

    def check_narrow_wrapping(self, slide, sn):
        """Detect text boxes where narrow width forces choppy, short-line wrapping.
        
        This catches the common visual bug where a text box is too narrow for its
        content, producing ugly 2-3 word lines that are hard to read.
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                box_w_in = emu_to_in(shape.width)
                box_h_in = emu_to_in(shape.height)
            except:
                continue

            tf = shape.text_frame
            total_text = tf.text.strip()
            if not total_text or len(total_text) < 15:
                continue

            font_size = _get_dominant_font_size(tf)

            # Estimate characters that fit per line in this box
            # Rule of thumb: at Xpt font, ~(box_width_inches * 72 / font_size * 1.7) chars
            # Aptos averages ~0.55 of the em-width per character
            chars_per_line = max(int(box_w_in * 72 / font_size * 1.7), 1)

            # Find the longest paragraph — that's the one most likely to wrap badly
            for para in tf.paragraphs:
                para_text = para.text.strip()
                if not para_text or len(para_text) < 15:
                    continue

                # How many lines will this paragraph wrap into?
                wrapped_lines = max(len(para_text) / chars_per_line, 1)

                if wrapped_lines >= 2:
                    avg_words_per_line = len(para_text.split()) / wrapped_lines

                    # Severe: less than ~3 words per line AND more than 3 lines
                    # This means text is breaking into choppy fragments
                    if avg_words_per_line < 3.0 and wrapped_lines > 3:
                        snippet = para_text[:60].replace('\n', ' ')
                        self.add('ERROR', sn, 'Narrow wrap',
                                 f'Text wraps into ~{int(wrapped_lines)} very short lines '
                                 f'(~{avg_words_per_line:.1f} words/line)',
                                 f'"{shape.name}" is {box_w_in:.1f}" wide at {font_size}pt — '
                                 f'text: "{snippet}..."')

                    # Moderate: less than 4 words per line with multiple wraps
                    elif avg_words_per_line < 4.0 and wrapped_lines > 2:
                        snippet = para_text[:60].replace('\n', ' ')
                        self.add('WARNING', sn, 'Narrow wrap',
                                 f'Text likely wraps into choppy lines '
                                 f'(~{avg_words_per_line:.1f} words/line across ~{int(wrapped_lines)} lines)',
                                 f'"{shape.name}" is {box_w_in:.1f}" wide at {font_size}pt — '
                                 f'widen the box, reduce text, or shrink font. '
                                 f'Text: "{snippet}..."')

                    # Mild: less than 5 words per line with many wraps — just a heads up
                    elif avg_words_per_line < 5.0 and wrapped_lines > 4 and chars_per_line < 30:
                        snippet = para_text[:60].replace('\n', ' ')
                        self.add('WARNING', sn, 'Narrow wrap',
                                 f'Narrow box may produce short, hard-to-read lines '
                                 f'(~{chars_per_line} chars/line)',
                                 f'"{shape.name}" is {box_w_in:.1f}" wide at {font_size}pt — '
                                 f'text: "{snippet}..."')

    def check_font_legibility(self, slide, sn):
        """Check that font sizes meet minimum legibility thresholds.
        
        Deloitte presentations are projected or shared on screen. Text below
        certain sizes becomes unreadable.
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                box_w_in = emu_to_in(shape.width)
                box_h_in = emu_to_in(shape.height)
            except:
                box_w_in = 1
                box_h_in = 1

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text.strip()
                    if not text:
                        continue
                    if not run.font.size:
                        continue

                    pt = run.font.size.pt
                    snippet = text[:40]

                    # Determine what kind of text this is by context
                    is_title = False
                    try:
                        pf = shape.placeholder_format
                        if pf and pf.type in (1, 3):
                            is_title = True
                    except (ValueError, AttributeError):
                        pass

                    # Check if this is a source/footnote line
                    is_source = (text.lower().startswith('source') or
                                 text.lower().startswith('note:') or
                                 text.lower().startswith('footnote'))

                    # Absolute floor: nothing below 6pt is ever readable on screen
                    if pt < 6:
                        self.add('ERROR', sn, 'Legibility',
                                 f'Font size {pt}pt is unreadable',
                                 f'"{snippet}" in "{shape.name}" — '
                                 f'minimum 7pt for footnotes/sources, 9pt for body text')

                    # Source citations below 7pt — enforce the 7pt floor
                    elif is_source and pt < 7:
                        self.add('WARNING', sn, 'Legibility',
                                 f'Source citation at {pt}pt below 7pt minimum',
                                 f'"{snippet}" in "{shape.name}" — '
                                 f'source lines should be 7-8pt for legibility when projected')

                    # Titles should be at least 14pt
                    elif is_title and pt < 14:
                        self.add('WARNING', sn, 'Legibility',
                                 f'Title text at {pt}pt may be too small',
                                 f'"{snippet}" — titles should be 14pt minimum, '
                                 f'recommended 18-24pt')

                    # Body text below 7pt is an error — never acceptable
                    elif not is_title and not is_source and pt < 7 and len(text) > 5:
                        self.add('ERROR', sn, 'Legibility',
                                 f'Body text at {pt}pt is too small',
                                 f'"{snippet}" in "{shape.name}" — '
                                 f'10pt minimum for body text, 7pt absolute floor')

                    # Body text at 8-9pt — warn, should be 10pt+
                    elif not is_title and not is_source and pt < 10 and len(text) > 15:
                        # Report once per shape to avoid flooding with per-run warnings
                        _key = f'{sn}:{shape.name}:body_below_10pt'
                        if not hasattr(self, '_body_warned'):
                            self._body_warned = set()
                        if _key not in self._body_warned:
                            self._body_warned.add(_key)
                            self.add('WARNING', sn, 'Legibility',
                                     f'Body text at {pt}pt below 10pt minimum',
                                     f'"{shape.name}" uses {pt}pt body text — '
                                     f'use 10-11pt for comfortable readability when projected. '
                                     f'If space-constrained, split across slides rather than shrinking font')

                    # Large font in tiny box — will wrap or clip
                    elif pt >= 18 and box_w_in < 2.0 and len(text) > 10:
                        self.add('WARNING', sn, 'Legibility',
                                 f'Large font ({pt}pt) in narrow box ({box_w_in:.1f}" wide)',
                                 f'"{snippet}" in "{shape.name}" — '
                                 f'text will wrap into very short lines')

    def check_zero_padded_numbers(self, slide, sn):
        """Flag zero-padded numbers in small shapes (likely numbered circles).
        
        Numbers like '01', '02' in oval/circle shapes will wrap to two lines
        because they're wider than single digits. Always use '1', '2', '3'.
        """
        import re
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Check for zero-padded numbers in small shapes (circles, ovals, small boxes)
            try:
                w = emu_to_in(shape.width)
                h = emu_to_in(shape.height)
            except:
                continue
            # Only check small shapes (< 1.2" in both dimensions) — these are circles/badges
            if w > 1.2 or h > 1.2:
                continue
            if re.match(r'^0\d$', text) and w < 0.70:
                self.add('WARNING', sn, 'Formatting',
                         f'Zero-padded "{text}" may not fit in {w:.2f}" circle',
                         f'"{shape.name}" ({w:.2f}" × {h:.2f}") — '
                         f'use "{int(text)}" instead of "{text}" to prevent '
                         f'text wrapping inside the circle')

    def check_card_whitespace(self, slide, sn):
        """Detect card/box rectangles where text fills less than 40% of the area.
        
        This catches the common anti-pattern where cards have static heights
        but short text content, leaving large empty colored blocks.
        Looks for filled rectangles (non-text) that have a text box overlapping
        them, and checks the text fill ratio.
        """
        shapes_data = []
        for shape in slide.shapes:
            try:
                left = emu_to_in(shape.left)
                top = emu_to_in(shape.top)
                w = emu_to_in(shape.width)
                h = emu_to_in(shape.height)
            except:
                continue
            has_text = shape.has_text_frame and shape.text_frame.text.strip()
            text = shape.text_frame.text.strip() if has_text else ''
            shapes_data.append({
                'shape': shape, 'name': shape.name,
                'left': left, 'top': top, 'width': w, 'height': h,
                'area': w * h, 'has_text': bool(has_text), 'text': text,
                'right': left + w, 'bottom': top + h
            })

        # Find filled rectangles (potential card bodies) > 1.5" tall and > 2" wide
        for rect in shapes_data:
            if rect['has_text'] or rect['height'] < 1.5 or rect['width'] < 2.0:
                continue
            if rect['area'] < 4.0:  # Skip small elements
                continue
            
            # Find text boxes that overlap this rectangle
            overlapping_text = []
            for tb in shapes_data:
                if not tb['has_text'] or tb is rect:
                    continue
                # Check if text box is mostly inside this rectangle
                if (tb['left'] >= rect['left'] - 0.2 and
                    tb['right'] <= rect['right'] + 0.2 and
                    tb['top'] >= rect['top'] - 0.2 and
                    tb['bottom'] <= rect['bottom'] + 0.2):
                    overlapping_text.append(tb)
            
            if not overlapping_text:
                continue
            
            # Estimate how much vertical space the text actually needs
            total_text = ' '.join(t['text'] for t in overlapping_text)
            if not total_text or len(total_text) < 10:
                continue
            
            # Rough estimate: at 10pt, ~12 chars/inch width, ~6 lines/inch height
            chars_per_line = int(rect['width'] * 11)  # conservative for 10pt
            est_lines = max(1, len(total_text) / max(chars_per_line, 1))
            est_text_height = est_lines * (10 / 72) * 1.4 + 0.3  # lines + padding
            
            fill_ratio = est_text_height / rect['height'] if rect['height'] > 0 else 1
            
            if fill_ratio < 0.40 and rect['height'] > 2.0:
                empty_pct = int((1 - fill_ratio) * 100)
                self.add('WARNING', sn, 'Card whitespace',
                         f'Card body is ~{empty_pct}% empty ({rect["height"]:.1f}" tall, '
                         f'text needs ~{est_text_height:.1f}")',
                         f'"{rect["name"]}" — use dynamic card height sizing. '
                         f'Set height to ~{est_text_height:.1f}" instead of {rect["height"]:.1f}" '
                         f'to eliminate the empty block')

    def _check_pale_green_fills(self, slide, sn):
        """Flag card-sized shapes with pale green fills on white-background slides.
        
        Pale green (F1F6E4) as a default card fill is the most common 'auto-generated'
        visual defect. Cards should match the slide background with an outline instead.
        """
        PALE_GREEN_RGB = 'F1F6E4'
        count = 0
        for shape in slide.shapes:
            try:
                w = emu_to_in(shape.width)
                h = emu_to_in(shape.height)
                if w < 1.5 or h < 0.5:
                    continue
                if shape.fill and shape.fill.type is not None:
                    fc = str(shape.fill.fore_color.rgb)
                    if fc == PALE_GREEN_RGB:
                        count += 1
            except:
                continue
        if count >= 3:
            self.add('WARNING', sn, 'Card fill',
                     f'{count} shapes use pale green (F1F6E4) fill',
                     f'Cards should match the slide background with an outline — '
                     f'use white fill + gray border on white slides, '
                     f'dark gray fill on dark slides. Use make_card() helper.')

    def check_text_density(self, slide, sn):
        """Check for too much text crammed into a small area.
        
        Even if text technically fits, high character density makes slides 
        look wall-of-text heavy and violates presentation best practices.
        """
        # Also check for pale green card fills on white-background slides
        self._check_pale_green_fills(slide, sn)
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                box_w_in = emu_to_in(shape.width)
                box_h_in = emu_to_in(shape.height)
            except:
                continue

            area = box_w_in * box_h_in
            if area < 0.3:
                continue

            total_text = shape.text_frame.text.strip()
            if not total_text:
                continue

            char_count = len(total_text)
            font_size = _get_dominant_font_size(shape.text_frame)

            # Characters per square inch — adjusted for font size
            # At 10pt, ~200 chars/sq-in is dense; at 8pt, ~300 is dense
            density = char_count / area
            density_threshold = 350 * (8 / max(font_size, 6))  # Scales with font size

            if density > density_threshold and char_count > 100:
                self.add('WARNING', sn, 'Density',
                         f'Very high text density ({int(density)} chars/sq-in)',
                         f'"{shape.name}": {char_count} characters in '
                         f'{box_w_in:.1f}" × {box_h_in:.1f}" at {font_size}pt — '
                         f'consider splitting across slides or using a larger box')

    def check_font_size_jumps(self, slide, sn):
        """Detect jarring font size changes within a single text frame.
        
        A shape that mixes e.g. 24pt and 9pt text (other than a clear 
        title+body pattern) usually looks like a formatting mistake.
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            sizes = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.text.strip():
                        sizes.append(run.font.size.pt)

            if len(sizes) < 2:
                continue

            unique = sorted(set(sizes))
            if len(unique) < 2:
                continue

            smallest = unique[0]
            largest = unique[-1]
            ratio = largest / smallest if smallest > 0 else 0

            # A 3:1+ ratio within the same shape is suspicious
            # (e.g. 8pt body with 28pt title in the same text box)
            if ratio >= 3.0 and largest - smallest > 12:
                self.add('WARNING', sn, 'Size jump',
                         f'Large font size range within one shape '
                         f'({smallest}pt to {largest}pt, {ratio:.1f}× ratio)',
                         f'"{shape.name}" — may look inconsistent. '
                         f'Consider separate shapes for title vs. body text')

    def check_divider_consistency(self):
        """Flag decks that mix different section divider designs.

        A polished deck uses ONE divider style throughout. Mixing
        photo-backed dividers with solid-color ones, or using 3+ different
        layout indices for dividers, looks unfinished.
        """
        # Divider layout indices: 13-22 in the Deloitte template
        divider_range = set(range(13, 23))
        divider_layouts = []

        for i, slide in enumerate(self.prs.slides):
            sn = i + 1
            try:
                layout_name = slide.slide_layout.name
            except:
                continue

            # Check if this slide uses a divider layout
            layout_idx = None
            for idx, layout in enumerate(self.prs.slide_layouts):
                if layout == slide.slide_layout:
                    layout_idx = idx
                    break

            if layout_idx is not None and layout_idx in divider_range:
                divider_layouts.append((sn, layout_idx, layout_name))

        if len(divider_layouts) >= 2:
            unique_layouts = set(idx for _, idx, _ in divider_layouts)
            if len(unique_layouts) > 1:
                layout_list = ', '.join(
                    f'slide {sn}: layout {idx} ({name})'
                    for sn, idx, name in divider_layouts
                )
                self.add('WARNING', None, 'Divider consistency',
                         f'Section dividers use {len(unique_layouts)} different layouts',
                         f'Pick ONE divider style and use it for all section breaks. '
                         f'Current: {layout_list}')

    def check_slide_density(self):
        """Flag slides with an excessive number of text-heavy shapes.
        
        Neuroscience of Winning data: dense slides create stress and 
        reduce attention. This catches "wall of text" slides.
        """
        for i, slide in enumerate(self.prs.slides):
            sn = i + 1
            text_shapes = 0
            total_chars = 0

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if text and len(text) > 10:
                    text_shapes += 1
                    total_chars += len(text)

            # More than ~1500 characters on a single slide is very dense
            if total_chars > 1500:
                self.add('WARNING', sn, 'Slide density',
                         f'Slide has ~{total_chars} characters of text',
                         f'Dense slides reduce audience engagement. '
                         f'Consider splitting content or using visuals.')

            # More than 18 separate text shapes suggests a cluttered layout.
            # Structured grids (2x3 cards, comparison tables, before/after layouts)
            # legitimately use 13-18 text elements. Only flag truly excessive counts.
            if text_shapes > 18:
                self.add('WARNING', sn, 'Slide density',
                         f'Slide has {text_shapes} text elements',
                         f'Complex layouts with many text boxes risk '
                         f'overlaps and hard-to-follow reading order')

    def check_content_quality(self):
        """Content quality heuristics — checks narrative and communication quality.
        
        These go beyond brand/formatting to catch common executive communication
        failures: weak titles, missing sources, unsupported claims, monotony.
        """
        import re

        # ── Descriptive title detection ──
        # Only flag generic single-word topic labels. Compound titles like
        # "Questions We Help You Answer" are fine for pitches and marketing.
        WEAK_TITLE_PATTERNS = [
            r'^(?:overview|background|introduction|agenda|summary|'
            r'analysis|discussion|update|status|appendix|'
            r'recommendations?|conclusion|highlights)$'
        ]
        weak_title_re = re.compile('|'.join(WEAK_TITLE_PATTERNS), re.IGNORECASE)
        
        consecutive_same_layout = 0
        prev_layout_name = None
        title_count = 0
        weak_title_count = 0

        for sn, slide in enumerate(self.prs.slides, 1):
            layout_name = slide.slide_layout.name

            # Skip title slides (deck openers), dividers, and end slides for content checks
            # "Title Only" is a content layout, not a title slide — don't skip it
            is_content = not any(kw in layout_name.lower() 
                               for kw in ['title slide', 'divider', 'end slide'])
            
            if is_content:
                # ── Check for descriptive (weak) titles ──
                for shape in slide.shapes:
                    try:
                        pf = shape.placeholder_format
                        if pf and pf.type in (1, 3):  # TITLE or CENTER_TITLE
                            title_text = shape.text_frame.text.strip()
                            if title_text:
                                title_count += 1
                                if weak_title_re.match(title_text):
                                    weak_title_count += 1
                                    self.add('WARNING', sn, 'Content quality',
                                             f'Descriptive title: "{title_text}"',
                                             f'Slide titles should express conclusions, not topics. '
                                             f'Rewrite as an assertion: what is the takeaway?')
                            break
                    except (ValueError, AttributeError):
                        continue

                # ── Check for unsourced quantitative claims ──
                has_numbers = False
                has_source = False
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    # Check for quantitative content (percentages, dollar amounts, large numbers)
                    if re.search(r'\d{2,}[%$€£¥]|[$€£¥]\d|(?:billion|million|bps|bbl)', text, re.I):
                        has_numbers = True
                    if text.lower().startswith('source') or 'source:' in text.lower():
                        has_source = True
                
                if has_numbers and not has_source:
                    self.add('WARNING', sn, 'Content quality',
                             'Quantitative claims without source footnote',
                             'Slide contains numbers/percentages but no source citation. '
                             'Add a source line or mark claims as "[source required]"')

            # ── Check for layout monotony (3+ consecutive same layout) ──
            if layout_name == prev_layout_name and is_content:
                consecutive_same_layout += 1
            else:
                consecutive_same_layout = 1
            prev_layout_name = layout_name

            if consecutive_same_layout >= 3:
                self.add('WARNING', sn, 'Content quality',
                         f'{consecutive_same_layout} consecutive slides use the same layout',
                         f'Layout "{layout_name}" repeated — vary the visual structure '
                         f'to maintain audience engagement')

        # ── Deck-level: too many weak titles ──
        if title_count > 0 and weak_title_count > title_count * 0.4:
            self.add('WARNING', None, 'Content quality',
                     f'{weak_title_count} of {title_count} content slides have '
                     f'descriptive (non-assertive) titles',
                     f'Rewrite titles as conclusions. The title sequence should tell '
                     f'the deck\'s story when read without body text')

    def check_contrast(self, slide, sn):
        """Basic contrast check — light text on light backgrounds or vice versa."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Try to get fill color
            bg_hex = None
            try:
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    bg_hex = rgb_to_hex(shape.fill.fore_color.rgb)
            except:
                pass

            if bg_hex is None:
                continue

            bg_lum = _luminance(bg_hex)

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    try:
                        if run.font.color and run.font.color.rgb:
                            txt_hex = rgb_to_hex(run.font.color.rgb)
                            txt_lum = _luminance(txt_hex)

                            # WCAG AA requires 4.5:1 for normal text
                            ratio = _contrast_ratio(bg_lum, txt_lum)
                            if ratio < 3.0:
                                snippet = run.text[:30].strip()
                                self.add('WARNING', sn, 'Contrast',
                                         f'Low contrast ({ratio:.1f}:1)',
                                         f'Text "#{txt_hex}" on bg "#{bg_hex}": '
                                         f'"{snippet}"')
                    except:
                        pass

    def check_consistency(self):
        """Cross-slide consistency checks."""
        # Check for mixed font sizes in titles across slides
        title_sizes = []
        for i, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                try:
                    pf = shape.placeholder_format
                    if pf is None:
                        continue
                    ptype = pf.type
                except (ValueError, AttributeError):
                    continue
                if ptype in (1, 3):  # TITLE or CENTER_TITLE
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                title_sizes.append((i+1, run.font.size.pt))
                                break
                        break

        if title_sizes:
            sizes = set(s for _, s in title_sizes)
            if len(sizes) > 2:
                self.add('WARNING', None, 'Consistency',
                         f'Title font sizes vary across slides: {sorted(sizes)}pt',
                         'Consider standardizing title sizes')

    # ── Report ───────────────────────────────────────────────────────────

    def report(self):
        """Print formatted report."""
        errors = [i for i in self.issues if i.level == 'ERROR']
        warnings = [i for i in self.issues if i.level == 'WARNING']
        infos = [i for i in self.issues if i.level == 'INFO']

        print()
        print("=" * 70)
        print(f"  DELOITTE PPTX VALIDATION REPORT")
        print(f"  File: {os.path.basename(self.path)}")
        print(f"  Slides: {len(self.prs.slides)}")
        print("=" * 70)

        if errors:
            print(f"\n❌ ERRORS ({len(errors)}):")
            for i in errors:
                print(str(i))

        if warnings:
            print(f"\n⚠️  WARNINGS ({len(warnings)}):")
            for i in warnings:
                print(str(i))

        if infos:
            print(f"\nℹ️  INFO ({len(infos)}):")
            for i in infos:
                print(str(i))

        print()
        print("-" * 70)
        if errors:
            print(f"  RESULT: ❌ FAIL — {len(errors)} error(s), {len(warnings)} warning(s)")
            print(f"  Fix the errors above before delivering.")
        elif warnings:
            print(f"  RESULT: ⚠️  PASS WITH WARNINGS — {len(warnings)} warning(s)")
            print(f"  Review warnings above; most may be acceptable.")
        else:
            print(f"  RESULT: ✅ PASS — No issues found")
        print("-" * 70)
        print()

        return len(errors) == 0

    def to_json(self):
        """Return issues as JSON for programmatic use."""
        return json.dumps([{
            'level': i.level,
            'slide': i.slide,
            'category': i.category,
            'message': i.message,
            'detail': i.detail,
        } for i in self.issues], indent=2)

# ── Text Layout Helpers ──────────────────────────────────────────────────────

def _get_dominant_font_size(text_frame):
    """Get the most common font size in a text frame (in pt). Defaults to 10."""
    sizes = {}
    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.font.size and run.text.strip():
                pt = run.font.size.pt
                weight = len(run.text)
                sizes[pt] = sizes.get(pt, 0) + weight
    if not sizes:
        return 10
    return max(sizes, key=sizes.get)

def _estimate_wrap(text, box_w_in, box_h_in, font_size_pt):
    """Estimate how text wraps in a box. Returns dict with line/capacity info.
    
    Uses Aptos-calibrated metrics: at Xpt, average char width ≈ X * 0.48 points,
    so chars_per_line ≈ (box_width_in * 72) / (font_size * 0.48)
    Line height ≈ font_size * 1.25 (with default spacing)
    """
    char_width_factor = 0.48  # Average Aptos char width as fraction of font size
    line_height_factor = 1.35  # Line height including spacing

    chars_per_line = max(int((box_w_in * 72) / (font_size_pt * char_width_factor)), 1)
    line_height_pt = font_size_pt * line_height_factor
    capacity_lines = max(int((box_h_in * 72) / line_height_pt), 1)

    # Count explicit newlines + estimate word-wrap lines
    paragraphs = text.split('\n')
    est_lines = 0
    for para in paragraphs:
        para = para.strip()
        if not para:
            est_lines += 1
        else:
            est_lines += max(len(para) / chars_per_line, 1)

    overflows = est_lines > capacity_lines * 1.2 and len(text) > 20

    return {
        'chars_per_line': chars_per_line,
        'est_lines': int(est_lines),
        'capacity': capacity_lines,
        'overflows': overflows,
    }

# ── Contrast Helpers ─────────────────────────────────────────────────────────

def _luminance(hex_color):
    """Relative luminance per WCAG 2.0."""
    r, g, b = int(hex_color[0:2], 16)/255, int(hex_color[2:4], 16)/255, int(hex_color[4:6], 16)/255
    r = r/12.92 if r <= 0.03928 else ((r+0.055)/1.055)**2.4
    g = g/12.92 if g <= 0.03928 else ((g+0.055)/1.055)**2.4
    b = b/12.92 if b <= 0.03928 else ((b+0.055)/1.055)**2.4
    return 0.2126*r + 0.7152*g + 0.0722*b

def _contrast_ratio(l1, l2):
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)

# ── CLI ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description='Validate a Deloitte-branded PPTX')
    parser.add_argument('pptx_file', help='Path to the .pptx file to validate')
    parser.add_argument('--strict', action='store_true',
                        help='Treat font warnings as errors')
    parser.add_argument('--json', action='store_true',
                        help='Output issues as JSON instead of formatted report')
    parser.add_argument('--visual', action='store_true',
                        help='Also generate slide images for visual inspection')
    args = parser.parse_args()

    if not os.path.exists(args.pptx_file):
        print(f"ERROR: File not found: {args.pptx_file}")
        sys.exit(1)

    v = Validator(args.pptx_file, strict=args.strict)
    v.check_all()

    if args.json:
        print(v.to_json())
    else:
        passed = v.report()

    if args.visual:
        print("Generating slide images for visual review...")
        pptx_dir = os.path.dirname(os.path.abspath(args.pptx_file))
        pptx_base = os.path.splitext(os.path.basename(args.pptx_file))[0]
        try:
            import subprocess
            # Convert to PDF first
            subprocess.run([
                'python', '/mnt/skills/public/pptx/scripts/office/soffice.py',
                '--headless', '--convert-to', 'pdf', args.pptx_file
            ], check=True, capture_output=True)
            pdf_path = os.path.splitext(args.pptx_file)[0] + '.pdf'
            if os.path.exists(pdf_path):
                subprocess.run([
                    'pdftoppm', '-jpeg', '-r', '150', pdf_path,
                    os.path.join(pptx_dir, f'{pptx_base}-slide')
                ], check=True, capture_output=True)
                print(f"  Slide images saved as {pptx_base}-slide-*.jpg")
            else:
                print("  Could not generate PDF for visual review")
        except Exception as e:
            print(f"  Visual review failed: {e}")

    sys.exit(0 if v.stats.get('ERROR', 0) == 0 else 1)


if __name__ == '__main__':
    main()
